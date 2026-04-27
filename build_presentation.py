"""
Build the US Air Attack Pilot Orientation PowerPoint.
Coulson Aviation design system — built from scratch, no template.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Design tokens ──
BG_COLOR = "1A1B1E"
TEXT_PRIMARY = RGBColor(0xE8, 0xEA, 0xED)
TEXT_SECONDARY = RGBColor(0x91, 0x96, 0xA0)
ACCENT_BLUE = RGBColor(0x5B, 0x7B, 0xD7)
BRAND_RED = RGBColor(0xD0, 0x00, 0x00)
STATUS_GREEN = RGBColor(0x5A, 0x9A, 0x6B)
STATUS_CAUTION = RGBColor(0xA8, 0x83, 0x2E)
STATUS_ERROR = RGBColor(0xC4, 0x58, 0x58)

FONT_BODY = "Inter"
FONT_MONO = "JetBrains Mono"

LOGO_PATH = r"C:\Users\jonas\Projects\Design\Coulson Templates\coulson-logo-stacked-white.png"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

OUTPUT = "US Air Attack Pilot Orientation.pptx"


# ── Helpers ──

def set_slide_bg(slide, hex_color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def remove_bullets(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    for child in list(pPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag.startswith("bu"):
            pPr.remove(child)
    etree.SubElement(pPr, qn("a:buNone"))


def set_para_spacing(paragraph, space_after_pt=12):
    pPr = paragraph._p.get_or_add_pPr()
    spc_aft = etree.SubElement(pPr, qn("a:spcAft"))
    spc_pts = etree.SubElement(spc_aft, qn("a:spcPts"))
    spc_pts.set("val", str(int(space_after_pt * 100)))


def add_text_box(slide, left, top, width, height, text, font_name=FONT_BODY,
                 font_size=Pt(18), font_color=TEXT_PRIMARY, bold=False,
                 alignment=PP_ALIGN.LEFT, all_caps=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    if all_caps:
        # python-pptx doesn't directly support caps, set via XML
        rPr = run._r.get_or_add_rPr()
        rPr.set("cap", "all")
    remove_bullets(p)
    return txBox, tf, p


def add_accent_line(slide, left, top, width=Inches(0.6), height=Pt(2),
                    color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_logo(slide, logo_path, left, top, width):
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, left, top, width=width)


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_section_label(slide, text):
    add_text_box(
        slide,
        left=Inches(0.8), top=Inches(0.5),
        width=Inches(4), height=Inches(0.3),
        text=text,
        font_size=Pt(9), font_color=ACCENT_BLUE,
        bold=True, all_caps=True,
    )


def add_slide_title(slide, text):
    add_text_box(
        slide,
        left=Inches(0.8), top=Inches(1.0),
        width=Inches(9), height=Inches(0.6),
        text=text,
        font_size=Pt(32), font_color=TEXT_PRIMARY, bold=True,
    )
    # Accent line below title
    add_accent_line(
        slide,
        left=Inches(0.8), top=Inches(1.65),
    )


def add_body_paragraph(tf, text, font_size=Pt(18), font_color=TEXT_SECONDARY,
                       bold=False, font_name=FONT_BODY, space_after=12):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    remove_bullets(p)
    set_para_spacing(p, space_after)
    return p


def add_content_logo(slide):
    """Small logo, bottom-right corner."""
    logo_w = Inches(0.7)
    logo_left = SLIDE_WIDTH - logo_w - Inches(0.5)
    logo_top = SLIDE_HEIGHT - Inches(0.9)
    add_logo(slide, LOGO_PATH, logo_left, logo_top, logo_w)


def new_slide(prs):
    """Add a blank slide and set the background."""
    layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    return slide


def build_content_slide(prs, section, title, lines, notes):
    """
    Build a standard content slide.
    lines: list of dicts with keys:
        text, type ('key'|'supporting'), and optional font_color override
    """
    slide = new_slide(prs)
    add_section_label(slide, section)
    add_slide_title(slide, title)

    # Body text frame
    body_top = Inches(2.0)
    txBox = slide.shapes.add_textbox(
        Inches(0.8), body_top, Inches(8), Inches(4.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for line in lines:
        text = line["text"]
        line_type = line.get("type", "supporting")
        color_override = line.get("color", None)

        if line_type == "key":
            sz = Pt(20)
            clr = TEXT_PRIMARY
            bld = True
        else:
            sz = Pt(18)
            clr = TEXT_SECONDARY
            bld = False

        if color_override:
            clr = color_override

        if first:
            # Use the existing first paragraph
            p = tf.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.name = FONT_BODY
            run.font.size = sz
            run.font.color.rgb = clr
            run.font.bold = bld
            remove_bullets(p)
            set_para_spacing(p, 12)
            first = False
        else:
            add_body_paragraph(tf, text, font_size=sz, font_color=clr, bold=bld)

    add_content_logo(slide)
    set_notes(slide, notes)
    return slide


# ── Build ──

def build_presentation():
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ================================================================
    # SLIDE 1 — TITLE SLIDE
    # ================================================================
    slide = new_slide(prs)

    # Title
    add_text_box(
        slide,
        left=Inches(1.5), top=Inches(2.4),
        width=Inches(8), height=Inches(1.0),
        text="Air Attack Pilot Orientation",
        font_size=Pt(44), font_color=TEXT_PRIMARY, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Brand-red accent line (centered)
    line_w = Inches(0.6)
    add_accent_line(
        slide,
        left=Inches(1.5) + (Inches(8) - line_w) // 2,
        top=Inches(3.5),
        color=BRAND_RED,
    )

    # Subtitle
    add_text_box(
        slide,
        left=Inches(1.5), top=Inches(3.8),
        width=Inches(8), height=Inches(0.5),
        text="COULSON AVIATION \u2014 US OPERATIONS",
        font_size=Pt(14), font_color=TEXT_SECONDARY,
        alignment=PP_ALIGN.CENTER, all_caps=True,
    )

    # Logo — right of center
    logo_w = Inches(1.5)
    add_logo(slide, LOGO_PATH,
             left=Inches(10.5), top=Inches(2.2), width=logo_w)

    set_notes(slide, (
        "Welcome everyone. This session is about getting aligned before you go "
        "out and start flying. You\u2019ve been through the technical ground school "
        "\u2014 this is about how Coulson does things, what we expect from you, and "
        "how to make the most of the air attack role."
    ))

    # ================================================================
    # SLIDE 2 — Why We're Here
    # ================================================================
    build_content_slide(prs, "OPENING", "Why We\u2019re Here", [
        {"text": "You\u2019re about to go out on your own", "type": "key"},
        {"text": "No captain to shadow \u2014 what you learn in this room is what you take to the field", "type": "supporting"},
    ], notes=(
        "Unlike our tanker crews, where a new pilot flies with an experienced "
        "captain, you\u2019re going out solo. There\u2019s no one in the left seat to show "
        "you the ropes. So we need to make sure you leave this room knowing how "
        "we operate, what we expect, and who to call when you need help."
    ))

    # ================================================================
    # SLIDE 3 — The Opportunity
    # ================================================================
    build_content_slide(prs, "OPENING", "The Opportunity", [
        {"text": "The best seat in the house to learn aerial firefighting", "type": "key"},
        {"text": "After a couple of seasons, you should be able to teach the class", "type": "supporting"},
        {"text": "Merit-based progression to air tanker", "type": "supporting"},
    ], notes=(
        "This is the message I really want to land. The air attack seat is the "
        "best seat in the house for learning aerial firefighting. You\u2019re overhead "
        "for four hours at a time. You\u2019ve got one of the most experienced ground "
        "firefighters in the country sitting next to you. You\u2019re watching fire "
        "behavior develop. You\u2019re hearing target descriptions. You\u2019re seeing how "
        "tactical decisions play out. When you move to the tanker, that\u2019s all "
        "gone \u2014 you\u2019re on scene for five minutes and then you\u2019re back to the "
        "reload base. Use this time. It\u2019s more valuable than most people realize "
        "until it\u2019s gone."
    ))

    # ================================================================
    # SLIDE 4 — What Coulson Expects
    # ================================================================
    build_content_slide(prs, "THE COULSON STANDARD", "What Coulson Expects", [
        {"text": "Professionalism \u2014 Initiative \u2014 Consistency", "type": "key"},
        {"text": "You are Coulson\u2019s face", "type": "supporting"},
    ], notes=(
        "You represent this company every time you walk onto a tanker base, "
        "every time you key the mic, every interaction you have with agency "
        "personnel. How you carry yourself matters. We\u2019ve built a strong "
        "reputation and we need you to uphold it."
    ))

    # ================================================================
    # SLIDE 5 — Tanker Base Conduct
    # ================================================================
    build_content_slide(prs, "THE COULSON STANDARD", "Tanker Base Conduct", [
        {"text": "You are a guest", "type": "key"},
        {"text": "Coulson polo \u00b7 Closed-toe shoes \u00b7 Clean up after yourself \u00b7 Build relationships", "type": "supporting"},
    ], notes=(
        "We are guests at every tanker base we operate from. Clean up after "
        "yourself. Be polite. Wear the Coulson polo and closed-toe shoes \u2014 "
        "present yourself in a way that reflects well on the company. Build "
        "relationships with the people you work around. These are the techs who "
        "keep your airplane flying, the base personnel who support your "
        "operation, and the agency staff who make it all happen. Treat them well."
    ))

    # ================================================================
    # SLIDE 6 — Your Support Structure
    # ================================================================
    build_content_slide(prs, "THE COULSON STANDARD", "Your Support Structure", [
        {"text": "You\u2019re flying alone, but you\u2019re not on your own", "type": "key"},
        {"text": "Open door \u2014 Chief Pilot, Chief of Training, Assistant Chief Pilot", "type": "supporting"},
        {"text": "Tell us early", "type": "supporting"},
    ], notes=(
        "Just because you\u2019re single pilot doesn\u2019t mean you\u2019re on your own. You "
        "can come to me, you can talk to our Chief of Training, you can talk to "
        "your Assistant Chief Pilot. We have an open door. If something is going "
        "on in the field \u2014 a difficult ATGS relationship, an operational issue, "
        "something you\u2019re unsure about \u2014 tell us. We\u2019d rather hear about it "
        "early and help you work through it than find out later when it\u2019s become "
        "a real problem. We\u2019re going to hold you accountable, but we\u2019re also "
        "going to have your back."
    ))

    # ================================================================
    # SLIDE 7 — Being Involved
    # ================================================================
    build_content_slide(prs, "THE COULSON STANDARD", "Being Involved", [
        {"text": "Don\u2019t just fly the plane \u2014 learn the mission", "type": "key"},
        {"text": "You are part of the aerial supervision team", "type": "supporting"},
    ], notes=(
        "This is probably the most important slide in this whole presentation. "
        "Your primary job is to fly the airplane and provide a stable platform. "
        "But if that\u2019s ALL you do, you\u2019re missing the point. You\u2019re overhead for "
        "four hours at a time. You\u2019ve got a front-row seat to watch fire "
        "behavior develop, to hear how tactical decisions are made, to see how "
        "retardant works on the fire. This is your education for the tanker. The "
        "pilots who get the most out of this role are the ones who are paying "
        "attention to all of it \u2014 not just flying the orbit."
    ))

    # ================================================================
    # SLIDE 8 — The ATGS Relationship
    # ================================================================
    build_content_slide(prs, "THE COULSON STANDARD", "The ATGS Relationship", [
        {"text": "Be the pilot they want to fly with", "type": "key"},
        {"text": "Be curious \u00b7 Ask questions \u00b7 Debrief every mission", "type": "supporting"},
        {"text": "The fire world is small", "type": "supporting"},
    ], notes=(
        "If you ask an ATGS who their favorite pilot is, they\u2019ll almost always "
        "tell you it\u2019s the one who was engaged. The one who wanted to learn. The "
        "one who asked questions during transit and at the debrief. They value "
        "that, and those relationships matter in this business. The fire world "
        "is small. Your reputation \u2014 good or bad \u2014 follows you. Be the pilot "
        "people want to work with."
    ))

    # ================================================================
    # SLIDE 9 — Daily Readiness
    # ================================================================
    build_content_slide(prs, "HOW WE OPERATE", "Daily Readiness", [
        {"text": "Ready at start time \u2014 not getting ready at start time", "type": "key"},
        {"text": "Weather awareness \u00b7 Gear staged \u00b7 Fuel planned", "type": "supporting"},
    ], notes=(
        "When the availability time hits, you need to be ready to go. That "
        "means showing up early enough to have the airplane preflighted, fueled, "
        "and ready. Not standing out on the ramp doing a walk-around when the "
        "phone rings. Stay aware of the weather, know what\u2019s going on in your "
        "area, and have your gear staged so that when dispatch calls, you can be "
        "airborne quickly."
    ))

    # ================================================================
    # SLIDE 10 — Dispatch and Mission Prep
    # ================================================================
    build_content_slide(prs, "HOW WE OPERATE", "Dispatch and Mission Prep", [
        {"text": "FRAT \u2014 a living document throughout the day", "type": "supporting"},
        {"text": "Pre-mission brief with AAS", "type": "supporting"},
        {"text": "Dispatch form + FRAT finalized = cleared for dispatch", "type": "key"},
    ], notes=(
        "The FRAT is a living document throughout the day. You start it during "
        "your walk-around, update it when conditions change or after dispatch, "
        "and finalize it at end of day. Before you launch on a mission, you "
        "brief with your AAS \u2014 what are the objectives, what frequencies are we "
        "on, how are we dividing the workload. You are not cleared for dispatch "
        "until both the dispatch form and FRAT are complete."
    ))

    # ================================================================
    # SLIDE 11 — Administrative Standards
    # ================================================================
    build_content_slide(prs, "HOW WE OPERATE", "Administrative Standards", [
        {"text": "EFL dailies \u2014 same evening, every evening", "type": "key"},
        {"text": "Standard email signature and format", "type": "supporting"},
        {"text": "Flight times match ATGS records \u2014 to the minute", "type": "supporting"},
        {"text": "All times in LOCAL", "type": "supporting"},
    ], notes=(
        "Admin consistency matters because you\u2019re all operating independently. "
        "We need everyone doing things the same way. EFL dailies go out the same "
        "evening \u2014 not the next morning, not two days later. Use the standard "
        "company email signature. The content in your daily email matters \u2014 "
        "include the relevant information. Your flight times need to match the "
        "ATGS records to the minute. If there\u2019s a discrepancy, it creates "
        "problems. And remember \u2014 all times are local, not UTC."
    ))

    # ================================================================
    # SLIDE 12 — Flight Following and Comms
    # ================================================================
    build_content_slide(prs, "HOW WE OPERATE", "Flight Following and Comms", [
        {"text": "AFF for all USFS flights", "type": "supporting"},
        {"text": "WhatsApp group for non-agency flights", "type": "supporting"},
    ], notes=(
        "For any flight under a USFS agreement, you\u2019ll be on AFF \u2014 Automated "
        "Flight Following. For non-agency flights \u2014 repositioning, ferry, "
        "anything not under an agency dispatch \u2014 use the WhatsApp group. Post "
        "your aircraft, where you\u2019re going, ETA, passengers, and fuel endurance. "
        "Someone needs to know where you are at all times."
    ))

    # ================================================================
    # SLIDE 13 — Crew Change and Handover
    # ================================================================
    build_content_slide(prs, "HOW WE OPERATE", "Crew Change and Handover", [
        {"text": "Handover form at every crew change", "type": "supporting"},
        {"text": "No surprises for the inbound crew", "type": "key"},
    ], notes=(
        "When you hand off to the next crew, the handover form needs to be "
        "complete and sent to the right people. Think about what the inbound "
        "pilot needs to know: aircraft status, any discrepancies, what\u2019s been "
        "going on operationally, anything unusual. They should be able to read "
        "your handover and step into the operation seamlessly."
    ))

    # ================================================================
    # SLIDE 14 — Citation Operations Field Guide
    # ================================================================
    build_content_slide(prs, "HOW WE OPERATE", "Citation Operations Field Guide", [
        {"text": "Day-to-day operational reference", "type": "supporting"},
        {"text": "Fuel \u00b7 OOS \u00b7 FMC \u00b7 Base ops \u00b7 Travel \u00b7 Contacts", "type": "supporting"},
    ], notes=(
        "The Citation Operations Field Guide covers the day-to-day procedures "
        "you\u2019ll need to reference \u2014 fuel cards, what to do if the airplane goes "
        "out of service, FMC database updates, expense reports, travel booking, "
        "and key contacts. It\u2019s there for you when you need it."
    ))

    # ================================================================
    # SLIDE 15 — The US Fire Season
    # ================================================================
    build_content_slide(prs, "US OPERATIONS CONTEXT", "The US Fire Season", [
        {"text": "Presenter: current season contracts and basing", "type": "supporting", "color": ACCENT_BLUE},
        {"text": "Exclusive use vs. call when needed", "type": "supporting"},
    ], notes=(
        "UPDATE THIS SLIDE EACH SEASON with current contract information.\n\n"
        "Walk through what contracts we have this year, where they are, whether "
        "they\u2019re exclusive use or call when needed, and what that means for the "
        "pilots\u2019 day-to-day. Talk about where they can expect to be based and "
        "the kinds of flights they may be asked to do outside their normal "
        "contract tasking \u2014 repositioning, ferry flights, maintenance runs. Set "
        "expectations: long days, possible diversions, and being ready to go "
        "wherever the work is."
    ))

    # ================================================================
    # SLIDE 16 — GACC Structure and Dispatch
    # ================================================================
    build_content_slide(prs, "US OPERATIONS CONTEXT", "GACC Structure and Dispatch", [
        {"text": "National \u2192 Geographic Area \u2192 Local Unit", "type": "key"},
        {"text": "How you get tasked", "type": "supporting"},
    ], notes=(
        "The US wildland fire dispatch system is organized around Geographic "
        "Area Coordination Centers \u2014 GACCs. There are several across the "
        "country, and they coordinate resource ordering within their areas. "
        "Above them is NIFC \u2014 the National Interagency Fire Center \u2014 which "
        "handles national-level coordination and resource movement between "
        "geographic areas. When things get busy, resources get moved around the "
        "country. Understanding this system helps you understand why you might "
        "get sent somewhere unexpected."
    ))

    # ================================================================
    # SLIDE 17 — Agency Relationships
    # ================================================================
    build_content_slide(prs, "US OPERATIONS CONTEXT", "Agency Relationships", [
        {"text": "USFS \u00b7 BLM \u00b7 CAL FIRE", "type": "supporting"},
        {"text": "Partners, not customers", "type": "key"},
    ], notes=(
        "You may be working under contracts with USFS, BLM, CAL FIRE, or other "
        "agencies. Each has their own way of doing things. The key thing to "
        "understand is that these are partnerships. Agency personnel \u2014 your "
        "ATGS, dispatch, tanker base managers \u2014 are people you work alongside. "
        "Treat them as partners. Build those relationships. It makes everything "
        "work better."
    ))

    # ================================================================
    # SLIDE 18 — Regulatory Framework
    # ================================================================
    build_content_slide(prs, "US OPERATIONS CONTEXT", "Regulatory Framework", [
        {"text": "FAR Part 135 \u00b7 Squawk 1255 \u00b7 NWCG PMS 505", "type": "supporting"},
        {"text": "TFRs \u2014 one incident\u2019s dispatch does not clear another\u2019s TFR", "type": "supporting"},
    ], notes=(
        "All of our flights under USFS agreements operate under Part 135 \u2014 no "
        "exceptions. Squawk 1255 for firefighting ops. One thing to be aware of "
        "when it gets busy: just because you\u2019ve been dispatched to an incident "
        "doesn\u2019t mean you can fly through another incident\u2019s TFR on the way "
        "there. Each TFR is its own airspace restriction. NWCG PMS 505 is the "
        "national standard for aerial supervision \u2014 it defines how all of this "
        "is supposed to work."
    ))

    # ================================================================
    # SLIDE 19 — Your Job
    # ================================================================
    build_content_slide(prs, "THE AIR ATTACK ROLE", "Your Job", [
        {"text": "Stable platform \u00b7 Separation \u00b7 Anticipate", "type": "supporting"},
        {"text": "Stay ahead of the mission", "type": "key"},
    ], notes=(
        "Your primary responsibility is flying the airplane and providing a "
        "stable platform so the AAS can do their job. Maintain separation, stay "
        "ahead of the mission, and manage the systems. Anticipate what the AAS "
        "is going to need \u2014 don\u2019t wait to be asked. If you can see that fuel is "
        "going to be a factor, say something early. If weather is building, flag "
        "it. Be proactive."
    ))

    # ================================================================
    # SLIDE 20 — What You Don't Do
    # ================================================================
    build_content_slide(prs, "THE AIR ATTACK ROLE", "What You Don\u2019t Do", [
        {"text": "No tactical radio transmissions \u2014 except safety of flight", "type": "key"},
        {"text": "Monitor \u00b7 Intercom \u00b7 Back up the AAS", "type": "supporting"},
    ], notes=(
        "You do not transmit on tactical frequencies unless it\u2019s a "
        "safety-of-flight issue. The AAS handles all of the tactical "
        "communication \u2014 target descriptions, airtanker coordination, "
        "everything. You listen, you monitor, and you coordinate with the AAS "
        "on intercom. The one exception: if the AAS is completely task-saturated "
        "\u2014 say they\u2019re heads down on FM with ground resources \u2014 and a critical "
        "call comes in on air-to-air, you can answer it and then brief the AAS "
        "on intercom."
    ))

    # ================================================================
    # SLIDE 21 — CRM in the Air Attack
    # ================================================================
    build_content_slide(prs, "THE AIR ATTACK ROLE", "CRM in the Air Attack", [
        {"text": "If something is developing, say it", "type": "key"},
        {"text": "Back each other up", "type": "supporting"},
    ], notes=(
        "Good CRM in the air attack means being a proactive partner. You\u2019re "
        "monitoring all the frequencies \u2014 if you hear traffic that the AAS might "
        "have missed, call it out on intercom. If fuel state is becoming a "
        "factor, don\u2019t wait until it\u2019s critical. If you see weather building or "
        "you\u2019ve got a system anomaly, say something. Withholding information \u2014 "
        "even unintentionally, by just not speaking up \u2014 erodes the CRM "
        "relationship. You back each other up."
    ))

    # ================================================================
    # SLIDE 22 — Key Tactical Reminders (SPECIAL: 2x2 grid)
    # ================================================================
    slide = new_slide(prs)
    add_section_label(slide, "THE AIR ATTACK ROLE")
    add_slide_title(slide, "Key Tactical Reminders")

    # 2x2 grid of data points using JetBrains Mono for numbers
    grid_data = [
        ("1,000 FT", "Vertical separation"),
        ("1,500 AGL", "Hard floor"),
        ("170 KIAS", "Within FTA"),
        ("150 KIAS", "Maneuvering"),
    ]
    grid_left = [Inches(0.8), Inches(4.5), Inches(0.8), Inches(4.5)]
    grid_top = [Inches(2.2), Inches(2.2), Inches(3.8), Inches(3.8)]

    for i, (num_text, label_text) in enumerate(grid_data):
        # Large number in JetBrains Mono, accent blue
        add_text_box(
            slide,
            left=grid_left[i], top=grid_top[i],
            width=Inches(3.5), height=Inches(0.6),
            text=num_text,
            font_name=FONT_MONO, font_size=Pt(28),
            font_color=ACCENT_BLUE, bold=True,
        )
        # Label below in secondary color
        add_text_box(
            slide,
            left=grid_left[i], top=grid_top[i] + Inches(0.55),
            width=Inches(3.5), height=Inches(0.4),
            text=label_text,
            font_name=FONT_BODY, font_size=Pt(14),
            font_color=TEXT_SECONDARY,
        )

    # Key line below the grid
    add_text_box(
        slide,
        left=Inches(0.8), top=Inches(5.2),
        width=Inches(8), height=Inches(0.5),
        text="Communicate \u2014 Clearance \u2014 Comply",
        font_size=Pt(20), font_color=TEXT_PRIMARY, bold=True,
    )

    add_content_logo(slide)
    set_notes(slide, (
        "You covered all of this in detail during the ground school. These are "
        "the key numbers to keep in your head. 1,000 feet vertical separation. "
        "1,500 AGL minimum orbit altitude. 170 knots or less in the FTA, 150 or "
        "less when maneuvering. Overhead pattern entry when practical. And never "
        "enter an FTA without clearance \u2014 communicate, get clearance, comply."
    ))

    # ================================================================
    # SLIDE 23 — Audio Panel Management
    # ================================================================
    build_content_slide(prs, "THE AIR ATTACK ROLE", "Audio Panel Management", [
        {"text": "Prioritize with volume \u2014 never turn a frequency off", "type": "key"},
        {"text": "Air-to-air highest \u2192 Dispatch lowest", "type": "supporting"},
        {"text": "You are the backup", "type": "supporting"},
    ], notes=(
        "When you\u2019re monitoring four frequencies at once, use the individual "
        "volume controls on the audio panel to manage them. Air-to-air is your "
        "highest priority \u2014 keep it loudest. Dispatch is probably your lowest "
        "priority \u2014 turn it down so it\u2019s not competing. But never turn a "
        "frequency off. You are the backup for the ATGS. Your collective role is "
        "managing the airspace and maintaining situational awareness. If you\u2019ve "
        "got a radio turned off, you can\u2019t hear what\u2019s going on, and you can\u2019t "
        "back them up."
    ))

    # ================================================================
    # SLIDE 24 — Know Your Limits
    # ================================================================
    build_content_slide(prs, "SAFETY & BRIGHT LINES", "Know Your Limits", [
        {"text": "Your envelope \u2014 not the airplane\u2019s", "type": "key"},
        {"text": "That mindset starts here", "type": "supporting"},
    ], notes=(
        "This is a direct message. The airplane has an operating envelope, and "
        "so do you. Yours is smaller. If a pilot is out of touch with their own "
        "ability \u2014 if they think they\u2019re hot shit and they\u2019re really not \u2014 that "
        "leads them to do things in an airplane that are beyond their skill "
        "level. And that\u2019s how people get killed in this business. Humility "
        "isn\u2019t weakness. It\u2019s what keeps you alive. The pilot who thinks they\u2019re "
        "beyond learning is the one who puts themselves in a situation they "
        "can\u2019t handle. That mindset starts here. Now. Be the pilot who knows "
        "what they don\u2019t know and keeps working to close the gap."
    ))

    # ================================================================
    # SLIDE 25 — Safety Culture
    # ================================================================
    build_content_slide(prs, "SAFETY & BRIGHT LINES", "Safety Culture", [
        {"text": "We back you \u2014 no penalty, no questions", "type": "key"},
        {"text": "Report through SMS", "type": "supporting"},
    ], notes=(
        "If you make a decision based on safety, we back you. No penalty, no "
        "questions. We would rather you make a conservative call and have a "
        "conversation about it than push into something that doesn\u2019t feel right. "
        "Report through SMS when something happens \u2014 near misses, safety "
        "concerns, anything. That\u2019s how we learn and how we get better as a "
        "company."
    ))

    # ================================================================
    # SLIDE 26 — Bright Lines (SPECIAL: accent-blue squares)
    # ================================================================
    slide = new_slide(prs)
    add_section_label(slide, "SAFETY & BRIGHT LINES")
    add_slide_title(slide, "Bright Lines")

    bright_items = [
        "Integrity",
        "Willful disregard for SOPs",
        "Representing Coulson poorly",
        "Conduct that undermines trust",
    ]

    item_top = Inches(2.1)
    for item in bright_items:
        # Small accent-blue square
        sq_size = Inches(0.12)
        sq = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), item_top + Inches(0.05),
            sq_size, sq_size,
        )
        sq.fill.solid()
        sq.fill.fore_color.rgb = ACCENT_BLUE
        sq.line.fill.background()

        # Text next to the square
        add_text_box(
            slide,
            left=Inches(1.15), top=item_top,
            width=Inches(7), height=Inches(0.35),
            text=item,
            font_size=Pt(18), font_color=TEXT_SECONDARY,
        )
        item_top += Inches(0.5)

    # Key line below
    add_text_box(
        slide,
        left=Inches(0.8), top=item_top + Inches(0.3),
        width=Inches(8), height=Inches(0.5),
        text="Mistakes we work through \u2014 this is about character",
        font_size=Pt(20), font_color=TEXT_PRIMARY, bold=True,
    )

    add_content_logo(slide)
    set_notes(slide, (
        "We\u2019ve done well hiring good people and I\u2019m not standing up here "
        "lecturing you about basic decency. But I do want to be clear about "
        "where the lines are. Integrity is non-negotiable. If something happened "
        "\u2014 you made a mistake, something went wrong \u2014 tell us. We\u2019ll work "
        "through it. Lying about it is worse than whatever the thing was. "
        "Willful disregard for SOPs, representing the company poorly, conduct "
        "that undermines the trust of the people you work with \u2014 those are "
        "bright lines. This isn\u2019t about operational mistakes. Everyone makes "
        "those and we\u2019ll work through them together. This is about character."
    ))

    # ================================================================
    # SLIDE 27 — FRAT and Tactical Pause (SPECIAL: colored circles)
    # ================================================================
    slide = new_slide(prs)
    add_section_label(slide, "SAFETY & BRIGHT LINES")
    add_slide_title(slide, "FRAT and Tactical Pause")

    frat_items = [
        (STATUS_GREEN, "Proceed"),
        (STATUS_CAUTION, "Mitigate"),
        (STATUS_ERROR, "Not authorized"),
    ]

    circle_top = Inches(2.2)
    for color, label in frat_items:
        # Colored circle (oval shape with equal w/h)
        circle_size = Inches(0.3)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), circle_top + Inches(0.02),
            circle_size, circle_size,
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()

        # Label text
        add_text_box(
            slide,
            left=Inches(1.3), top=circle_top,
            width=Inches(4), height=Inches(0.35),
            text=label,
            font_size=Pt(18), font_color=TEXT_SECONDARY,
        )
        circle_top += Inches(0.55)

    # Supporting text
    add_text_box(
        slide,
        left=Inches(0.8), top=Inches(4.2),
        width=Inches(8), height=Inches(0.4),
        text="Either crew member can call a tactical pause",
        font_size=Pt(18), font_color=TEXT_SECONDARY,
    )

    # Key line
    add_text_box(
        slide,
        left=Inches(0.8), top=Inches(4.8),
        width=Inches(8), height=Inches(0.5),
        text="We\u2019d rather have a conversation than an incident",
        font_size=Pt(20), font_color=TEXT_PRIMARY, bold=True,
    )

    add_content_logo(slide)
    set_notes(slide, (
        "The FRAT gives you a structured way to assess risk. Green means go. "
        "Yellow means you need to identify and apply mitigations before "
        "proceeding. Red means the mission is not authorized. And at any point "
        "during a mission, either crew member can call a tactical pause \u2014 no "
        "penalty, no judgment. If something doesn\u2019t feel right, say so. We "
        "would rather have a conversation about it than deal with an incident."
    ))

    # ================================================================
    # SLIDE 28 — Training Form (TRN-04)
    # ================================================================
    build_content_slide(prs, "WHAT YOU\u2019RE BEING CHECKED ON", "Training Form (TRN-04)", [
        {"text": "What you\u2019ll be signed off on before your OPC", "type": "key"},
    ], notes=(
        "The training form covers everything you need to be signed off on "
        "before your OPC. It\u2019s organized into ground school competencies and "
        "flight training competencies. Your training pilot will work through "
        "this with you and sign off each area as you demonstrate proficiency."
    ))

    # ================================================================
    # SLIDE 29 — OPC Checkride (OPC-04)
    # ================================================================
    build_content_slide(prs, "WHAT YOU\u2019RE BEING CHECKED ON", "OPC Checkride (OPC-04)", [
        {"text": "~1 hour evaluation flight \u2014 check pilot acts as AAS", "type": "supporting"},
        {"text": "Annual recurrency", "type": "supporting"},
    ], notes=(
        "The OPC is about a one-hour evaluation flight where the check pilot "
        "flies in the right seat acting as the AAS. You\u2019ll be evaluated on 26 "
        "items across six competency areas \u2014 everything from mission knowledge "
        "and FTA procedures to non-technical skills like CRM and "
        "decision-making. There are automatic failure gates: a safety event, a "
        "repeated deficiency, or checklist discipline failures will result in a "
        "Not Approved outcome. The OPC is required annually to maintain your "
        "qualification."
    ))

    # ================================================================
    # SLIDE 30 — CLOSING TITLE SLIDE
    # ================================================================
    slide = new_slide(prs)

    # Title
    add_text_box(
        slide,
        left=Inches(1.5), top=Inches(2.4),
        width=Inches(8), height=Inches(1.0),
        text="Chief Pilot\u2019s Message",
        font_size=Pt(44), font_color=TEXT_PRIMARY, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Brand-red accent line (centered)
    line_w = Inches(0.6)
    add_accent_line(
        slide,
        left=Inches(1.5) + (Inches(8) - line_w) // 2,
        top=Inches(3.5),
        color=BRAND_RED,
    )

    # Subtitle
    add_text_box(
        slide,
        left=Inches(1.5), top=Inches(3.8),
        width=Inches(8), height=Inches(0.5),
        text="QUESTIONS",
        font_size=Pt(14), font_color=TEXT_SECONDARY,
        alignment=PP_ALIGN.CENTER, all_caps=True,
    )

    # Logo — centered, larger
    logo_w = Inches(1.5)
    logo_left = Inches(1.5) + (Inches(8) - logo_w) // 2
    add_logo(slide, LOGO_PATH, logo_left, Inches(4.8), logo_w)

    set_notes(slide, (
        "You\u2019re part of something special at Coulson. We\u2019re investing in you "
        "and we take that seriously. The air attack seat is the best seat in the "
        "house to learn aerial firefighting \u2014 use it for what it\u2019s worth. "
        "Perform well, learn constantly, represent Coulson with pride. The "
        "tanker seat is earned, not given. This is where you earn it.\n\n"
        "Open the floor for questions."
    ))

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_presentation()
